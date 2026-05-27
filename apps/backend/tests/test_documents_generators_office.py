"""Tests de los generadores PPTX + XLSX + PDF (Fase 4.2).

Cada test re-abre el archivo generado con su lib (python-pptx,
openpyxl, pypdf vía pdfplumber) para verificar que es estructuralmente
válido. Sin MS Office.
"""

from __future__ import annotations

import io
from datetime import UTC, datetime

import pdfplumber
from openpyxl import load_workbook
from pptx import Presentation

from sqa_kb.documents.generators.base import GeneratedFile
from sqa_kb.documents.generators.pdf import PdfGenerator
from sqa_kb.documents.generators.pptx import PptxGenerator
from sqa_kb.documents.generators.xlsx import XlsxGenerator
from sqa_kb.documents.models import DocumentContent, QaPair


def _content(**overrides) -> DocumentContent:  # type: ignore[no-untyped-def]
    base = {
        "document_id": "PRES-roadmap-qa-2026-05-26",
        "title": "Roadmap de QA 2026",
        "category": "EST",
        "document_type": "PRES",
        "version": "2.0",
        "fecha": datetime(2026, 5, 26, tzinfo=UTC),
        "author_name": "Tester",
        "author_role": "QA Lead",
        "topic": "estrategia de automatización para el próximo año",
        "body_blocks": ("Pilar 1: shift-left.", "Pilar 2: contract testing."),
        "qa_pairs": (
            QaPair(question="¿Qué herramienta de e2e?", answer="Playwright."),
        ),
        "is_anonymized": False,
    }
    base.update(overrides)
    return DocumentContent(**base)  # type: ignore[arg-type]


# ===========================================================================
# PptxGenerator
# ===========================================================================


def test_pptx_returns_valid_presentation() -> None:
    out = PptxGenerator().generate(_content())
    assert isinstance(out, GeneratedFile)
    assert out.filename == "PRES-roadmap-qa-2026-05-26.pptx"
    assert out.media_type.endswith("presentationml.presentation")
    prs = Presentation(io.BytesIO(out.data))
    assert prs is not None


def test_pptx_slide_count() -> None:
    """1 portada + 1 tema + 2 contenido + 1 precision = 5 slides."""
    out = PptxGenerator().generate(_content())
    prs = Presentation(io.BytesIO(out.data))
    assert len(prs.slides) == 5


def test_pptx_cover_has_title() -> None:
    out = PptxGenerator().generate(_content())
    prs = Presentation(io.BytesIO(out.data))
    cover = prs.slides[0]
    texts = [
        shape.text_frame.text
        for shape in cover.shapes
        if shape.has_text_frame
    ]
    assert any("Roadmap de QA 2026" in t for t in texts)


def test_pptx_single_content_block_no_numbering() -> None:
    out = PptxGenerator().generate(_content(body_blocks=("Único bloque.",)))
    prs = Presentation(io.BytesIO(out.data))
    all_text = " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Contenido" in all_text
    # No debería numerar si hay un solo bloque.
    assert "Contenido 1" not in all_text


def test_pptx_omits_content_slides_when_no_blocks() -> None:
    out = PptxGenerator().generate(_content(body_blocks=(), qa_pairs=()))
    prs = Presentation(io.BytesIO(out.data))
    # Portada + tema = 2 slides.
    assert len(prs.slides) == 2


# ===========================================================================
# XlsxGenerator
# ===========================================================================


def test_xlsx_returns_valid_workbook() -> None:
    out = XlsxGenerator().generate(_content())
    assert out.filename == "PRES-roadmap-qa-2026-05-26.xlsx"
    assert out.media_type.endswith("spreadsheetml.sheet")
    wb = load_workbook(io.BytesIO(out.data))
    assert wb is not None


def test_xlsx_has_three_sheets_when_qa_present() -> None:
    out = XlsxGenerator().generate(_content())
    wb = load_workbook(io.BytesIO(out.data))
    assert wb.sheetnames == ["Metadata", "Contenido", "Precisiones"]


def test_xlsx_omits_precisiones_sheet_when_no_qa() -> None:
    out = XlsxGenerator().generate(_content(qa_pairs=()))
    wb = load_workbook(io.BytesIO(out.data))
    assert "Precisiones" not in wb.sheetnames
    assert wb.sheetnames == ["Metadata", "Contenido"]


def test_xlsx_metadata_values() -> None:
    out = XlsxGenerator().generate(_content())
    wb = load_workbook(io.BytesIO(out.data))
    meta = wb["Metadata"]
    # Construir dict propiedad→valor (saltando header en row 1).
    props = {meta.cell(r, 1).value: meta.cell(r, 2).value for r in range(2, meta.max_row + 1)}
    assert props["Título"] == "Roadmap de QA 2026"
    assert "EST" in props["Carpeta"]
    assert "PRES" in props["Tipo"]
    assert props["Versión"] == "2.0"
    assert props["Anonimizado"] == "No"


def test_xlsx_content_sheet_has_topic_and_blocks() -> None:
    out = XlsxGenerator().generate(_content())
    wb = load_workbook(io.BytesIO(out.data))
    content_ws = wb["Contenido"]
    col_b = [content_ws.cell(r, 2).value for r in range(2, content_ws.max_row + 1)]
    assert "estrategia de automatización para el próximo año" in col_b
    assert "Pilar 1: shift-left." in col_b
    assert "Pilar 2: contract testing." in col_b


def test_xlsx_anonymized_yes() -> None:
    out = XlsxGenerator().generate(_content(is_anonymized=True))
    wb = load_workbook(io.BytesIO(out.data))
    meta = wb["Metadata"]
    props = {meta.cell(r, 1).value: meta.cell(r, 2).value for r in range(2, meta.max_row + 1)}
    assert props["Anonimizado"] == "Sí"


# ===========================================================================
# PdfGenerator
# ===========================================================================


def test_pdf_returns_valid_pdf() -> None:
    out = PdfGenerator().generate(_content())
    assert out.filename == "PRES-roadmap-qa-2026-05-26.pdf"
    assert out.media_type == "application/pdf"
    # Magic bytes de PDF.
    assert out.data[:5] == b"%PDF-"


def test_pdf_extractable_text_has_title_and_sections() -> None:
    out = PdfGenerator().generate(_content())
    with pdfplumber.open(io.BytesIO(out.data)) as pdf:
        text = "\n".join(page.extract_text() or "" for page in pdf.pages)
    assert "Roadmap de QA 2026" in text
    assert "Tema" in text
    assert "Contenido" in text
    assert "Precisiones" in text
    assert "Playwright." in text


def test_pdf_escapes_special_chars_without_crashing() -> None:
    """Texto con `<`, `>`, `&` no debe romper el parser de reportlab."""
    out = PdfGenerator().generate(
        _content(
            topic="comparar a < b && c > d en el pipeline",
            body_blocks=("usar <tag> & validar > 0",),
        )
    )
    assert out.data[:5] == b"%PDF-"
    with pdfplumber.open(io.BytesIO(out.data)) as pdf:
        text = "\n".join(page.extract_text() or "" for page in pdf.pages)
    assert "a < b" in text or "a &lt; b" in text or "comparar" in text


def test_pdf_footer_has_page_number() -> None:
    out = PdfGenerator().generate(_content())
    with pdfplumber.open(io.BytesIO(out.data)) as pdf:
        text = "\n".join(page.extract_text() or "" for page in pdf.pages)
    assert "Página 1" in text
    assert "SQA" in text


def test_pdf_omits_optional_sections_when_empty() -> None:
    out = PdfGenerator().generate(_content(body_blocks=(), qa_pairs=()))
    with pdfplumber.open(io.BytesIO(out.data)) as pdf:
        text = "\n".join(page.extract_text() or "" for page in pdf.pages)
    assert "Tema" in text
    assert "Contenido" not in text
    assert "Precisiones" not in text
