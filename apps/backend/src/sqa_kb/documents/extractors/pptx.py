"""PptxExtractor — lee texto + estructura de un `.pptx` (Fase 4.3).

Usa python-pptx. Cada slide se modela como una `ExtractedSection`: el
título del slide (primer text frame, o "Slide N" si no hay) es el
`title`, y el resto de los text frames es el `content`.
"""

from __future__ import annotations

import io

from pptx import Presentation

from sqa_kb.documents.extractors.base import ExtractedDocument, ExtractedSection

EXTENSIONS = ("pptx",)


class PptxExtractor:
    """Implementa `DocumentExtractor` para `.pptx`."""

    @property
    def extensions(self) -> tuple[str, ...]:
        return EXTENSIONS

    def extract(self, data: bytes) -> ExtractedDocument:
        prs = Presentation(io.BytesIO(data))
        sections: list[ExtractedSection] = []
        all_text: list[str] = []

        for idx, slide in enumerate(prs.slides, start=1):
            frames = [
                shape.text_frame.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            if not frames:
                continue
            title = frames[0]
            body = "\n".join(frames[1:]).strip()
            sections.append(ExtractedSection(title=title or f"Slide {idx}", content=body))
            all_text.extend(frames)

        return ExtractedDocument(
            text="\n".join(all_text).strip(),
            sections=tuple(sections),
            page_count=len(prs.slides),
        )
