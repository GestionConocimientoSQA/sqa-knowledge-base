"""PptxGenerator — presentaciones PowerPoint con branding SQA (Fase 4.2).

Pensado para el tipo PRES, pero acepta cualquier `DocumentContent`.

Estructura de slides:
1. Portada: banda azul corp + título + tipo/carpeta + barra naranja.
2. Tema: un slide con el topic.
3. Contenido: un slide por cada `body_block` (título "Contenido N").
4. Precisiones: un slide por cada `QaPair`.

Branding: fondo blanco, títulos en azul corp (Exo 2), acentos naranja,
footer con `SQA · Gestión del Conocimiento`. Sin plantilla `.pptx`
binaria — todo programático.
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from sqa_kb.documents import branding
from sqa_kb.documents.branding import RgbColor
from sqa_kb.documents.doc_types import category_label, doc_type_label
from sqa_kb.documents.generators.base import GeneratedFile
from sqa_kb.documents.models import DocumentContent

MEDIA_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
EXTENSION = "pptx"

# Dimensiones de slide 16:9 (default de python-pptx en EMU).
_SLIDE_W = Emu(9144000)
_SLIDE_H = Emu(6858000)


def _rgb(color: RgbColor) -> RGBColor:
    return RGBColor.from_string(color.hex)


class PptxGenerator:
    """Genera `.pptx` con branding SQA. Implementa `DocumentGenerator`."""

    @property
    def extension(self) -> str:
        return EXTENSION

    @property
    def media_type(self) -> str:
        return MEDIA_TYPE

    def generate(self, content: DocumentContent) -> GeneratedFile:
        prs = Presentation()
        prs.slide_width = _SLIDE_W
        prs.slide_height = _SLIDE_H

        self._add_cover(prs, content)
        self._add_text_slide(prs, "Tema", [content.topic])
        if content.body_blocks:
            for i, block in enumerate(content.body_blocks, start=1):
                title = "Contenido" if len(content.body_blocks) == 1 else f"Contenido {i}"
                self._add_text_slide(prs, title, [block])
        if content.qa_pairs:
            for qa in content.qa_pairs:
                self._add_text_slide(prs, qa.question, [qa.answer])

        buffer = io.BytesIO()
        prs.save(buffer)
        return GeneratedFile(
            filename=f"{content.document_id}.{EXTENSION}",
            media_type=MEDIA_TYPE,
            data=buffer.getvalue(),
        )

    # -- slides -------------------------------------------------------------

    def _blank_slide(self, prs: Presentation):  # type: ignore[no-untyped-def]
        """Slide con layout en blanco (último layout = 'Blank' en el
        template default de python-pptx)."""
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _add_cover(self, prs: Presentation, content: DocumentContent) -> None:  # type: ignore[no-untyped-def]
        slide = self._blank_slide(prs)

        # Caption: tipo · carpeta.
        tipo = doc_type_label(content.document_type)
        carpeta = category_label(content.category)
        caption = self._add_textbox(
            slide, Emu(685800), Emu(1600200), Emu(7772400), Emu(457200)
        )
        self._set_text(
            caption,
            f"{tipo}  ·  {carpeta}".upper(),
            size=branding.SIZE_CAPTION + 2,
            color=branding.COLOR_TEXTO_SECUNDARIO,
            font=branding.FONT_BODY,
            bold=False,
        )

        # Título.
        title_box = self._add_textbox(
            slide, Emu(685800), Emu(2057400), Emu(7772400), Emu(1828800)
        )
        self._set_text(
            title_box,
            content.title,
            size=branding.SIZE_PORTADA_TITULO,
            color=branding.COLOR_TITULO,
            font=branding.FONT_DISPLAY,
            bold=True,
        )

        # Barra naranja debajo del título (firma SQA).
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Emu(685800),
            Emu(3962400),
            Emu(2743200),
            Emu(91440),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(branding.COLOR_ACENTO)
        bar.line.fill.background()

        self._add_footer(slide)

    def _add_text_slide(self, prs: Presentation, title: str, bodies: list[str]) -> None:  # type: ignore[no-untyped-def]
        slide = self._blank_slide(prs)

        title_box = self._add_textbox(
            slide, Emu(685800), Emu(457200), Emu(7772400), Emu(914400)
        )
        self._set_text(
            title_box,
            title,
            size=branding.SIZE_H1,
            color=branding.COLOR_SUBTITULO,
            font=branding.FONT_DISPLAY,
            bold=True,
        )

        body_box = self._add_textbox(
            slide, Emu(685800), Emu(1600200), Emu(7772400), Emu(4343400)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, body in enumerate(bodies):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = para.add_run()
            run.text = body
            run.font.name = branding.FONT_BODY
            run.font.size = Pt(branding.SIZE_BODY + 3)
            run.font.color.rgb = _rgb(branding.COLOR_CUERPO)

        self._add_footer(slide)

    # -- helpers ------------------------------------------------------------

    def _add_textbox(self, slide, left, top, width, height):  # type: ignore[no-untyped-def]
        box = slide.shapes.add_textbox(left, top, width, height)
        box.text_frame.word_wrap = True
        return box

    def _set_text(  # type: ignore[no-untyped-def]
        self, box, text: str, *, size: int, color: RgbColor, font: str, bold: bool
    ) -> None:
        tf = box.text_frame
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)

    def _add_footer(self, slide) -> None:  # type: ignore[no-untyped-def]
        footer = self._add_textbox(
            slide, Emu(685800), Emu(6400800), Emu(7772400), Emu(304800)
        )
        self._set_text(
            footer,
            branding.DOC_FOOTER_LEFT,
            size=branding.SIZE_FOOTER,
            color=branding.COLOR_TEXTO_SECUNDARIO,
            font=branding.FONT_BODY,
            bold=False,
        )
